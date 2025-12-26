# # make_modules/pptx_maker.py



"""
PowerPoint presentation generator for the Journal Lookup Tool.
Creates professional presentations with paper summaries and statistics.
"""



import os
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Tuple, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PresentationBuilder:
    """Builder class for creating PowerPoint presentations."""
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize presentation builder.
        
        Args:
            template_path: Optional path to PowerPoint template file
        """
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()
        
        # Layout indices for different slide types
        self.layouts = {
            'title': 0,          # Title slide
            'title_content': 1,  # Title and content
            'section_header': 2, # Section header
            'two_content': 3,    # Two content
            'comparison': 4,     # Comparison
            'title_only': 5,     # Title only
            'blank': 6,          # Blank
            'content_caption': 7, # Content with caption
            'picture_caption': 8  # Picture with caption
        }
    
    def create_text_slide(
        self,
        layout_type: str,
        slide_title: str,
        paragraphs: List[str],
        left: float = 1.0,
        top: float = 1.5,
        width: float = 8.0,
        height: float = 4.0,
        font_size: int = 12,
        font_name: str = "Calibri"
    ) -> None:
        """
        Create a slide with text content.
        
        Args:
            layout_type: Type of layout to use
            slide_title: Title for the slide
            paragraphs: List of paragraph texts
            left, top, width, height: Textbox dimensions in inches
            font_size: Font size in points
            font_name: Font family name
        """
        layout_index = self.layouts.get(layout_type, 5)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[layout_index])
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            self._format_title(slide.shapes.title)
        
        # Add textbox with content
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        for para_text in paragraphs:
            if not para_text.strip():
                # Add empty paragraph for spacing
                text_frame.add_paragraph()
                continue
                
            paragraph = text_frame.add_paragraph()
            paragraph.text = str(para_text)
            
            # Format paragraph
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    
    def create_paper_slide(self, component: Dict[str, Any]) -> None:
        """
        Create a slide for an individual paper.
        
        Args:
            component: Paper component dictionary
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])  # Blank layout
        
        # Create main content area
        left = top = Inches(0.2)
        width = self.presentation.slide_width - Inches(0.4)
        height = self.presentation.slide_height - Inches(0.4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        # Add paper title
        self._add_paper_title(text_frame, component)
        
        # Add spacing
        text_frame.add_paragraph()
        
        # Add abstract
        self._add_paper_abstract(text_frame, component)
        
        # Add spacing
        text_frame.add_paragraph()
        
        # Add publication details
        self._add_paper_details(text_frame, component)
        
        # Add authors
        self._add_paper_authors(text_frame, component)
        
        # Add institutions
        self._add_paper_institutions(text_frame, component)
    




    



    def create_statistics_slide(
        self,
        title: str,
        stats_data: Dict[str, int],
        max_items: int = 15
    ) -> None:
        """
        Create a slide with statistics.
        
        Args:
            title: Slide title
            stats_data: Dictionary of statistics
            max_items: Maximum number of items to display
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = title
        self._format_title(slide.shapes.title)
        
        # Sort and limit items
        sorted_items = sorted(stats_data.items(), key=lambda x: x[1], reverse=True)
        display_items = sorted_items[:max_items]
        
        if not display_items:
            # Use textbox for "no data" message
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = textbox.text_frame
            text_frame.clear()
            paragraph = text_frame.add_paragraph()
            paragraph.text = "No data available"
            paragraph.font.size = Pt(14)
            paragraph.font.italic = True
            return
        
        # Try to use content placeholder, fall back to textbox
        try:
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for name, count in display_items:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = f"{name}: {count}"  # Removed bullet point
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = "Calibri"
                
                # Add summary if there are more items
                if len(sorted_items) > max_items:
                    remaining = len(sorted_items) - max_items
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = f"... and {remaining} more"
                    paragraph.font.size = Pt(10)
                    paragraph.font.italic = True
            else:
                raise IndexError("No content placeholder found")
                
        except (IndexError, AttributeError):
            # Fall back to textbox
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = textbox.text_frame
            text_frame.clear()
            
            for name, count in display_items:
                paragraph = text_frame.add_paragraph()
                paragraph.text = f"• {name}: {count}"  # Keep bullet for textbox
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
            
            # Add summary if there are more items
            if len(sorted_items) > max_items:
                remaining = len(sorted_items) - max_items
                paragraph = text_frame.add_paragraph()
                paragraph.text = f"... and {remaining} more"
                paragraph.font.size = Pt(10)
                paragraph.font.italic = True







    
    def create_keywords_slide(self, components: List[Dict[str, Any]]) -> None:
        """
        Create a slide showing novel keywords found in papers.
        
        Args:
            components: List of paper components
        """
        # Extract papers with keywords
        papers_with_keywords = {}
        for component in components:
            keywords = component.get("Keywords", [])
            if (keywords and 
                keywords != ["No keywords available"] and 
                "No keywords" not in str(keywords)):
                title = component.get("Title", "Unknown paper")
                papers_with_keywords[title] = keywords
        
        if not papers_with_keywords:
            return  # Skip if no keywords found
        
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide.shapes.title.text = "Novel Associated Keywords Found"
        self._format_title(slide.shapes.title)
        
        # Create textbox for keywords
        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        for title, keywords in papers_with_keywords.items():
            # Add paper title
            title_para = text_frame.add_paragraph()
            title_para.text = f"• {self._truncate_text(title, 80)}:"
            title_para.font.size = Pt(11)
            title_para.font.bold = True
            title_para.font.name = "Calibri"
            
            # Add keywords with line wrapping
            keywords_text = ", ".join(keywords)
            wrapped_keywords = self._wrap_text(keywords_text, 85)
            
            for line in wrapped_keywords:
                keyword_para = text_frame.add_paragraph()
                keyword_para.text = f"   {line}"
                keyword_para.font.size = Pt(9)
                keyword_para.font.name = "Calibri"
                keyword_para.font.color.rgb = RGBColor(68, 68, 68)
            
            # Add spacing between papers
            text_frame.add_paragraph()
    







    def _format_title(self, title_shape) -> None:
        """Format slide title consistently."""
        if not title_shape or not title_shape.text_frame:
            return
            
        text_frame = title_shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.name = "Calibri"
                run.font.bold = True
                run.font.color.rgb = RGBColor(68, 114, 196)  # Professional blue
    
    def _add_paper_title(self, text_frame, component: Dict[str, Any]) -> None:
        """Add formatted paper title to text frame."""
        title = component.get("Title", "No title available")
        
        # Clean HTML tags
        clean_title = self._clean_html_tags(title)
        
        # Wrap title if too long
        wrapped_lines = self._wrap_text(clean_title, 75)
        
        for line in wrapped_lines:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(20)
            paragraph.font.name = "Calibri"
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(68, 114, 196)
    
    def _add_paper_abstract(self, text_frame, component: Dict[str, Any]) -> None:
        """Add formatted paper abstract to text frame."""
        abstract = component.get("Abstract", "No abstract available")
        
        if abstract == "No abstract available":
            paragraph = text_frame.add_paragraph()
            paragraph.text = "Abstract: Not available"
            paragraph.font.size = Pt(14)
            paragraph.font.italic = True
            return
        
        # Clean HTML tags
        clean_abstract = self._clean_html_tags(abstract)
        
        # Wrap abstract text
        wrapped_lines = self._wrap_text(clean_abstract, 115)
        
        for line in wrapped_lines:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Calibri"
    
    def _add_paper_details(self, text_frame, component: Dict[str, Any]) -> None:
        """Add paper publication details to text frame."""
        paragraph = text_frame.add_paragraph()
        
        # Journal
        journal = component.get("Journal", "Unknown journal")
        paragraph.text = f"Published in: {journal}\n"
        
        # Date
        date_info = component.get("Date", "Unknown date")
        formatted_date = self._format_date(date_info)
        paragraph.text += f"Publication Date: {formatted_date}\n"
        
        # DOI/Link
        link = self._extract_link(component.get("Link"))
        paragraph.text += f"DOI: {link}"
        
        paragraph.font.size = Pt(12)
        paragraph.font.name = "Calibri"
        paragraph.font.color.rgb = RGBColor(68, 68, 68)
    
    def _add_paper_authors(self, text_frame, component: Dict[str, Any]) -> None:
        """Add formatted authors to text frame."""
        authors = component.get("Authors", [])
        
        if not authors or authors == ["No authors available"]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = "Authors: Not available"
            paragraph.font.size = Pt(10)
            paragraph.font.italic = True
            return
        
        authors_text = "Authors: " + ", ".join(authors)
        wrapped_lines = self._wrap_text(authors_text, 140)
        
        for line in wrapped_lines:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(10)
            paragraph.font.name = "Calibri"
    
    def _add_paper_institutions(self, text_frame, component: Dict[str, Any]) -> None:
        """Add formatted institutions to text frame."""
        institutions = component.get("Institution", [])
        
        if not institutions or institutions == ["No institution listed"]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = "Institutions: Not available"
            paragraph.font.size = Pt(10)
            paragraph.font.italic = True
            return
        
        # Count institution frequency and show top ones
        institution_counts = {}
        for inst in institutions:
            institution_counts[inst] = institution_counts.get(inst, 0) + 1
        
        # Get top 3 institutions
        top_institutions = sorted(
            institution_counts.items(), 
            key=lambda x: x[1], 
            reverse=True
        )[:3]
        
        institution_text = "Institutions: " + "; ".join([
            f"{inst} ({count})" if count > 1 else inst 
            for inst, count in top_institutions
        ])
        
        wrapped_lines = self._wrap_text(institution_text, 140)
        
        for line in wrapped_lines:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(10)
            paragraph.font.name = "Calibri"
    
    def _clean_html_tags(self, text: str) -> str:
        """Remove HTML tags from text."""
        if not text:
            return ""
        
        # Remove common HTML tags
        cleaned = re.sub(r'</?(?:sub|sup|i|b|em|strong)>', '', text)
        return cleaned.strip()
    
    def _wrap_text(self, text: str, max_chars: int) -> List[str]:
        """Wrap text to specified character width."""
        if len(text) <= max_chars:
            return [text]
        
        lines = []
        while len(text) > max_chars:
            # Find last space before max_chars
            split_index = text.rfind(" ", 0, max_chars)
            if split_index == -1:
                split_index = max_chars
            
            lines.append(text[:split_index])
            text = text[split_index:].lstrip()
        
        if text:
            lines.append(text)
        
        return lines
    
    def _truncate_text(self, text: str, max_chars: int) -> str:
        """Truncate text to specified length with ellipsis."""
        if len(text) <= max_chars:
            return text
        return text[:max_chars-3] + "..."
    
    def _format_date(self, date_info: Any) -> str:
        """Format date information from various sources."""
        if isinstance(date_info, str):
            return date_info
        elif isinstance(date_info, list) and date_info:
            date_dict = date_info[0]
            if isinstance(date_dict, dict):
                year = date_dict.get("Year", "")
                month = date_dict.get("Month", "")
                day = date_dict.get("Day", "")
                if year and month and day:
                    return f"{year}/{month}/{day}"
                elif year:
                    return str(year)
        elif isinstance(date_info, dict):
            year = date_info.get("Year", "")
            month = date_info.get("Month", "")
            day = date_info.get("Day", "")
            if year and month and day:
                return f"{year}/{month}/{day}"
            elif year:
                return str(year)
        
        return "Not available"
    
    def _extract_link(self, link_info: Any) -> str:
        """Extract link/DOI from various formats."""
        if isinstance(link_info, str):
            return link_info
        elif isinstance(link_info, list) and link_info:
            return str(link_info[0])
        elif hasattr(link_info, 'value'):
            return str(link_info.value)
        else:
            return "Not available"
    
    def save(self, filename: str) -> None:
        """Save presentation to file."""
        self.presentation.save(filename)


def create_presentation(
    start_end_date: Tuple[str, str],
    config_file_dict: Dict[str, Any],
    components_keyword: List[Dict[str, Any]],
    components_orcid: List[Dict[str, Any]],
    keyword_frequency_dict: Dict[str, int],
    auto_mode: bool = False,
    version: str = "3.5.0",
    output_path: str = "publications.pptx",
    template_path: Optional[str] = None
) -> None:
    """
    Create a PowerPoint presentation with journal lookup results.
    
    Args:
        start_end_date: Date range tuple
        config_file_dict: Configuration dictionary
        components_keyword: Keyword-based search results
        components_orcid: ORCID-based search results
        keyword_frequency_dict: Keyword frequency statistics
        auto_mode: Whether running in automatic mode
        version: Tool version
        output_path: Output file path
        template_path: Optional PowerPoint template path
    """
    # Check if file exists and get permission
    if Path(output_path).exists() and not auto_mode:
        response = input(f"{output_path} already exists. Overwrite? (y/n): ").strip().lower()
        if response not in ['y', 'yes']:
            print("PowerPoint creation cancelled.")
            return
    
    # Initialize presentation builder
    builder = PresentationBuilder(template_path)
    
    # Extract data
    start_date, end_date = start_end_date
    email = config_file_dict.get('email', 'Not specified')
    journals = config_file_dict.get('journals', [])
    topics = config_file_dict.get('topics', [])
    authors = config_file_dict.get('authors', [])
    
    # Format end date for display
    end_date_display = end_date if end_date != '3000/01/01' else datetime.now().strftime('%Y/%m/%d')
    
    print("📊 Creating PowerPoint presentation...")
    
    # Slide 1: Title slide
    builder.create_text_slide(
        layout_type='title_only',
        slide_title=f"Journal Club Publication Watch Tool",
        paragraphs=[
            f"Watching Pubmed + CrossRef",
            f"Version {version}",
            "",
            f"For support, contact: aaronledray@gmail.com"
        ],
        font_size=16
    )
    


# -------------


    # Slide 2: Query information (basic)
    query_info = [
        f"User: {email}",
        "",
        f"Search period: {start_date} to {end_date_display}",
        "",
        f"Total keyword results: {len(components_keyword)}",
        f"Total author results: {len(components_orcid)}",
        f"Combined unique papers: {len(components_keyword) + len(components_orcid)}"
    ]
    
    if journals:
        query_info.extend(["", "Monitored journals:"])
        query_info.extend(journals[:10])  # Limit to first 10
        if len(journals) > 10:
            query_info.append(f"... and {len(journals) - 10} more")
    
    builder.create_text_slide(
        layout_type='title_only',
        slide_title="Search Summary",
        paragraphs=query_info,
        font_size=12
    )
    




    # Slide 3: Keywords overview
    if topics:
        keywords_info = ["Search keywords:"]
        
        # Format keywords with line wrapping
        keywords_text = ", ".join(topics)
        wrapped_keywords = builder._wrap_text(keywords_text, 80)
        keywords_info.extend(wrapped_keywords)
        
        builder.create_text_slide(
            layout_type='title_only',
            slide_title="Keywords Searched",
            paragraphs=keywords_info,
            font_size=12
        )
    




# -------------






    # Slide 4: Journal statistics (from keyword results)
    if components_keyword:
        journal_stats = {}
        for component in components_keyword:
            journal = component.get("Journal", "Unknown")
            if journal != "No journal available":
                journal_stats[journal] = journal_stats.get(journal, 0) + 1
        
        builder.create_statistics_slide(
            title="Journals Found (Keyword Search)",
            stats_data=journal_stats,
            max_items=15
        )
    


    # Slide 5: Keyword hit statistics
    if keyword_frequency_dict:
        filtered_keywords = {k: v for k, v in keyword_frequency_dict.items() if v > 0}
        
        builder.create_statistics_slide(
            title="Keyword Hit Statistics",
            stats_data=filtered_keywords,
            max_items=20
        )
    





# -------------





    # Slide 6: Novel keywords found
    if components_keyword:
        builder.create_keywords_slide(components_keyword)
    





# -------------







    # Individual paper slides (keyword-based)
    if components_keyword:
        for i, component in enumerate(components_keyword):
            print(f"   Adding keyword paper {i+1}/{len(components_keyword)}")
            builder.create_paper_slide(component)
    
    # Separator slide for ORCID results
    if components_orcid:
        builder.create_text_slide(
            layout_type='title_content',
            slide_title="Author-Based Results",
            paragraphs=[
                "The following papers were found through author/ORCID searches.",
                "",
                f"Total papers: {len(components_orcid)}",
                "",
                "These represent publications by researchers you are tracking."
            ],
            font_size=14
        )
        
        # Individual paper slides (ORCID-based)
        for i, component in enumerate(components_orcid):
            print(f"   Adding ORCID paper {i+1}/{len(components_orcid)}")
            builder.create_paper_slide(component)
    
    # Save presentation
    builder.save(output_path)
    print(f"  PowerPoint presentation saved: {output_path}")


# Legacy function name for backward compatibility
def get_pptx(
    start_end_date: Tuple[str, str],
    config_file_dict: Dict[str, Any],
    components_keyword: List[Dict[str, Any]],
    components_orcid: List[Dict[str, Any]],
    keyword_frequency_dict: Dict[str, int],
    auto_mode: bool = False,
    version: str = "3.5.0",
    **kwargs
) -> None:
    """Legacy wrapper for create_presentation function."""
    pptx_name = kwargs.get('pptx_name', 'publications.pptx')
    
    create_presentation(
        start_end_date=start_end_date,
        config_file_dict=config_file_dict,
        components_keyword=components_keyword,
        components_orcid=components_orcid,
        keyword_frequency_dict=keyword_frequency_dict,
        auto_mode=auto_mode,
        version=version,
        output_path=pptx_name
    )


if __name__ == "__main__":
    # Test the presentation builder
    print("Testing PowerPoint generator...")
    
    # Sample data
    sample_config = {
        'email': 'researcher@university.edu',
        'journals': ['Nature', 'Science', 'Cell'],
        'topics': ['machine learning', 'CRISPR'],
        'authors': ['0000-0000-0000-0001']
    }
    
    sample_keyword_components = [
        {
            'Title': 'Machine Learning in Drug Discovery',
            'Authors': ['John Doe', 'Jane Smith'],
            'Journal': 'Nature Biotechnology',
            'Date': [{'Year': '2024', 'Month': '01', 'Day': '15'}],
            'Abstract': 'This paper explores machine learning applications in drug discovery.',
            'Keywords': ['machine learning', 'drug discovery', 'artificial intelligence'],
            'Link': '10.1038/s41587-024-0001-1'
        }
    ]
    
    sample_orcid_components = [
        {
            'Title': 'CRISPR Gene Editing Advances',
            'Authors': ['Alice Wilson'],
            'Journal': 'Cell',
            'Date': [{'Year': '2024', 'Month': '02', 'Day': '20'}],
            'Abstract': 'Recent advances in CRISPR-Cas9 gene editing technology.',
            'Institution': ['Stanford University', 'MIT'],
            'Link': 'https://doi.org/10.1016/j.cell.2024.001'
        }
    ]
    
    # Test presentation creation
    create_presentation(
        start_end_date=('2024/01/01', '2024/12/31'),
        config_file_dict=sample_config,
        components_keyword=sample_keyword_components,
        components_orcid=sample_orcid_components,
        keyword_frequency_dict={'machine learning': 1, 'CRISPR': 1},
        auto_mode=True,
        version="3.5.0",
        output_path="test_presentation.pptx"
    )
    
    print("PowerPoint generator test completed.")